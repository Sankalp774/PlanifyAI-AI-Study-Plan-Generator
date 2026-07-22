"""Multi-format document ingestion."""

from __future__ import annotations

from pathlib import Path


def extract_text(path: str | Path) -> str:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".txt" or suffix == ".md":
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        import fitz

        text = []
        with fitz.open(path) as doc:
            for page in doc:
                text.append(page.get_text())
        return "\n".join(text)
    if suffix == ".docx":
        import docx

        d = docx.Document(path)
        return "\n".join(p.text for p in d.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return "\n".join(chunks)
    raise ValueError(f"Unsupported file type: {suffix}")
